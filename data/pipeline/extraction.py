"""
Content Extraction Module

Extracts raw content from various file formats:
- Text files (txt, md, csv, json)
- Documents (PDF, DOCX, PPTX)
- Web content (HTML, XML)
- Images (JPEG, PNG, TIFF, etc.)
- Audio (MP3, WAV, FLAC)
- Video (MP4, AVI, MKV)
- Code files
- Archives (ZIP)
"""

import os
import io
import json
import csv
import logging
from pathlib import Path
from typing import List, Dict, Optional, Any, Union, Tuple
from dataclasses import dataclass, field
from enum import Enum

import numpy as np

from .config import PipelineConfig, TokenizationPipelineConfig
from .ingestion import IngestedItem
from .preprocessing import PreprocessedItem

logger = logging.getLogger(__name__)


class ContentType(Enum):
    """Types of extracted content."""
    TEXT = "text"
    IMAGE_EMBEDDING = "image_embedding"
    AUDIO_EMBEDDING = "audio_embedding"
    VIDEO_EMBEDDING = "video_embedding"
    STRUCTURED_DATA = "structured_data"


@dataclass
class ExtractedContent:
    """Container for extracted content from any source."""
    # Source info
    source_path: str
    source_type: str
    
    # Text content
    text: Optional[str] = None
    
    # Embeddings for non-text modalities
    image_embeddings: Optional[np.ndarray] = None  # Shape: [num_patches, embed_dim]
    audio_embeddings: Optional[np.ndarray] = None  # Shape: [num_frames, embed_dim]
    video_embeddings: Optional[np.ndarray] = None  # Shape: [num_frames, embed_dim]
    
    # Structured data
    structured_data: Optional[Dict[str, Any]] = None
    
    # Metadata
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    # Content types present
    content_types: List[ContentType] = field(default_factory=list)
    
    # Processing info
    extraction_success: bool = True
    error_message: Optional[str] = None
    
    def has_text(self) -> bool:
        return self.text is not None and len(self.text) > 0
    
    def has_embeddings(self) -> bool:
        return any([
            self.image_embeddings is not None,
            self.audio_embeddings is not None,
            self.video_embeddings is not None,
        ])
    
    def total_embedding_tokens(self) -> int:
        """Calculate total number of embedding tokens."""
        total = 0
        if self.image_embeddings is not None:
            total += self.image_embeddings.shape[0]
        if self.audio_embeddings is not None:
            total += self.audio_embeddings.shape[0]
        if self.video_embeddings is not None:
            total += self.video_embeddings.shape[0]
        return total


class TextExtractor:
    """Extract text from various text-based formats."""
    
    @staticmethod
    def extract_txt(content: Union[str, bytes]) -> str:
        """Extract from plain text file."""
        if isinstance(content, bytes):
            return content.decode("utf-8", errors="ignore")
        return content
    
    @staticmethod
    def extract_json(content: Union[str, bytes]) -> Tuple[str, Dict]:
        """Extract from JSON file."""
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")
        
        try:
            data = json.loads(content)
            
            # Convert to text representation
            if isinstance(data, dict):
                # Extract text fields
                text_parts = []
                for key, value in data.items():
                    if isinstance(value, str):
                        text_parts.append(f"{key}: {value}")
                    elif isinstance(value, list):
                        text_parts.append(f"{key}: {', '.join(str(v) for v in value)}")
                text = "\n".join(text_parts)
                return text, data
            elif isinstance(data, list):
                # Handle JSONL-style data
                text_parts = []
                for item in data:
                    if isinstance(item, dict) and "text" in item:
                        text_parts.append(item["text"])
                    elif isinstance(item, str):
                        text_parts.append(item)
                return "\n".join(text_parts), {"items": data}
            else:
                return str(data), {"value": data}
        except json.JSONDecodeError:
            return content, {}
    
    @staticmethod
    def extract_csv(content: Union[str, bytes]) -> Tuple[str, List[Dict]]:
        """Extract from CSV file."""
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")
        
        rows = []
        text_parts = []
        
        try:
            reader = csv.DictReader(io.StringIO(content))
            for row in reader:
                rows.append(row)
                # Convert row to natural language
                row_text = ", ".join(f"{k}: {v}" for k, v in row.items() if v)
                text_parts.append(row_text)
        except Exception:
            # Fall back to raw text
            return content, []
        
        return "\n".join(text_parts), rows
    
    @staticmethod
    def extract_markdown(content: Union[str, bytes]) -> str:
        """Extract from Markdown file."""
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")
        
        # Basic markdown to text conversion
        import re
        
        # Remove code blocks
        text = re.sub(r"```[\s\S]*?```", "[CODE]", content)
        
        # Remove inline code
        text = re.sub(r"`[^`]+`", "[CODE]", text)
        
        # Remove headers markers but keep text
        text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
        
        # Remove link formatting but keep text
        text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
        
        # Remove bold/italic markers
        text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
        text = re.sub(r"\*([^*]+)\*", r"\1", text)
        text = re.sub(r"__([^_]+)__", r"\1", text)
        text = re.sub(r"_([^_]+)_", r"\1", text)
        
        return text


class DocumentExtractor:
    """Extract text from document formats (PDF, DOCX, PPTX)."""
    
    @staticmethod
    def extract_pdf(content: bytes, max_pages: int = 100) -> Tuple[str, Dict]:
        """
        Extract text from PDF using pdfplumber.
        
        pdfplumber provides superior text extraction with:
        - Better layout preservation
        - Table extraction support
        - Character-level positioning
        - Image extraction capabilities
        """
        metadata = {"pages": 0, "extraction_method": "none", "tables": 0}
        
        try:
            import pdfplumber
            
            text_parts = []
            table_count = 0
            
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                metadata["pages"] = len(pdf.pages)
                metadata["extraction_method"] = "pdfplumber"
                
                for i, page in enumerate(pdf.pages[:max_pages]):
                    try:
                        # Extract main text
                        text = page.extract_text()
                        if text:
                            text_parts.append(f"[Page {i+1}]\n{text}")
                        
                        # Extract tables if present
                        tables = page.extract_tables()
                        if tables:
                            table_count += len(tables)
                            for table in tables:
                                if table:
                                    # Convert table to text representation
                                    table_text = "\n".join(
                                        " | ".join(str(cell) if cell else "" for cell in row)
                                        for row in table if row
                                    )
                                    if table_text.strip():
                                        text_parts.append(f"[Table]\n{table_text}")
                    except Exception:
                        continue
            
            metadata["tables"] = table_count
            return "\n\n".join(text_parts), metadata
            
        except ImportError:
            logger.warning("pdfplumber not installed, run: pip install pdfplumber")
            return "", metadata
        except Exception as e:
            logger.error(f"Error extracting PDF: {e}")
            return "", metadata
    
    @staticmethod
    def extract_docx(content: bytes) -> Tuple[str, Dict]:
        """Extract text from DOCX."""
        metadata = {"paragraphs": 0}
        
        try:
            from docx import Document
            
            doc = Document(io.BytesIO(content))
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            metadata["paragraphs"] = len(text_parts)
            return "\n\n".join(text_parts), metadata
        except ImportError:
            logger.warning("python-docx not installed")
            return "", metadata
        except Exception as e:
            logger.error(f"Error extracting DOCX: {e}")
            return "", metadata
    
    @staticmethod
    def extract_pptx(content: bytes) -> Tuple[str, Dict]:
        """Extract text from PPTX."""
        metadata = {"slides": 0}
        
        try:
            from pptx import Presentation
            
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                
                if slide_texts:
                    text_parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
            
            metadata["slides"] = len(prs.slides)
            return "\n\n".join(text_parts), metadata
        except ImportError:
            logger.warning("python-pptx not installed")
            return "", metadata
        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}")
            return "", metadata


class WebExtractor:
    """Extract content from web formats (HTML, XML)."""
    
    @staticmethod
    def extract_html(content: Union[str, bytes]) -> Tuple[str, Dict]:
        """Extract text from HTML."""
        metadata = {"title": None, "links": 0}
        
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")
        
        try:
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(content, "html.parser")
            
            # Extract title
            title_tag = soup.find("title")
            if title_tag:
                metadata["title"] = title_tag.get_text().strip()
            
            # Remove script and style elements
            for element in soup(["script", "style", "nav", "footer", "header"]):
                element.decompose()
            
            # Get text
            text = soup.get_text(separator="\n", strip=True)
            
            # Count links
            metadata["links"] = len(soup.find_all("a"))
            
            return text, metadata
        except ImportError:
            # Fallback: basic regex extraction
            import re
            text = re.sub(r"<[^>]+>", " ", content)
            text = re.sub(r"\s+", " ", text)
            return text.strip(), metadata
        except Exception as e:
            logger.error(f"Error extracting HTML: {e}")
            return content if isinstance(content, str) else "", metadata
    
    @staticmethod
    def extract_xml(content: Union[str, bytes]) -> Tuple[str, Dict]:
        """Extract text from XML."""
        metadata = {"root_tag": None, "elements": 0}
        
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")
        
        try:
            import xml.etree.ElementTree as ET
            
            root = ET.fromstring(content)
            metadata["root_tag"] = root.tag
            
            text_parts = []
            element_count = 0
            
            def extract_text(element, depth=0):
                nonlocal element_count
                element_count += 1
                
                texts = []
                if element.text and element.text.strip():
                    texts.append(element.text.strip())
                
                for child in element:
                    texts.extend(extract_text(child, depth + 1))
                
                if element.tail and element.tail.strip():
                    texts.append(element.tail.strip())
                
                return texts
            
            text_parts = extract_text(root)
            metadata["elements"] = element_count
            
            return "\n".join(text_parts), metadata
        except Exception as e:
            logger.error(f"Error extracting XML: {e}")
            return "", metadata


class ImageExtractor:
    """Extract embeddings from images."""
    
    def __init__(self, config: TokenizationPipelineConfig):
        self.config = config
        self.encoder = None
        self._init_encoder()
    
    def _init_encoder(self):
        """Initialize image encoder."""
        try:
            if self.config.image_encoder == "clip":
                from transformers import CLIPProcessor, CLIPVisionModel
                self.processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
                self.encoder = CLIPVisionModel.from_pretrained("openai/clip-vit-base-patch32")
                self.encoder_type = "clip"
            else:
                # Default: simple patch-based encoding
                self.encoder_type = "patch"
            logger.info(f"Image encoder initialized: {self.encoder_type}")
        except ImportError:
            logger.warning("transformers not installed. Using patch-based encoding.")
            self.encoder_type = "patch"
        except Exception as e:
            logger.warning(f"Could not load image encoder: {e}")
            self.encoder_type = "patch"
    
    def extract(self, content: bytes) -> Tuple[Optional[np.ndarray], Dict]:
        """Extract embeddings from image."""
        metadata = {"width": 0, "height": 0, "channels": 0}
        
        try:
            from PIL import Image
            
            image = Image.open(io.BytesIO(content)).convert("RGB")
            metadata["width"], metadata["height"] = image.size
            metadata["channels"] = 3
            
            # Resize
            image = image.resize(self.config.image_size)
            
            if self.encoder_type == "clip" and self.encoder is not None:
                return self._extract_clip(image, metadata)
            else:
                return self._extract_patches(image, metadata)
        except ImportError:
            logger.warning("PIL not installed")
            return None, metadata
        except Exception as e:
            logger.error(f"Error extracting image: {e}")
            return None, metadata
    
    def _extract_clip(self, image, metadata: Dict) -> Tuple[np.ndarray, Dict]:
        """Extract using CLIP encoder."""
        inputs = self.processor(images=image, return_tensors="pt")
        outputs = self.encoder(**inputs)
        
        # Get patch embeddings (excluding CLS token)
        embeddings = outputs.last_hidden_state[:, 1:, :].detach().numpy()
        metadata["encoder"] = "clip"
        metadata["num_patches"] = embeddings.shape[1]
        
        return embeddings.squeeze(0), metadata
    
    def _extract_patches(self, image, metadata: Dict) -> Tuple[np.ndarray, Dict]:
        """Extract simple patch embeddings."""
        import numpy as np
        
        img_array = np.array(image)
        h, w, c = img_array.shape
        patch_size = self.config.image_patch_size
        
        # Extract patches
        patches = []
        for i in range(0, h, patch_size):
            for j in range(0, w, patch_size):
                patch = img_array[i:i+patch_size, j:j+patch_size, :]
                if patch.shape[:2] == (patch_size, patch_size):
                    # Flatten patch to embedding
                    patches.append(patch.flatten())
        
        if patches:
            embeddings = np.stack(patches)
            # Normalize
            embeddings = embeddings.astype(np.float32) / 255.0
            metadata["encoder"] = "patch"
            metadata["num_patches"] = len(patches)
            return embeddings, metadata
        
        return None, metadata


class AudioExtractor:
    """Extract embeddings from audio."""
    
    def __init__(self, config: TokenizationPipelineConfig):
        self.config = config
        self.encoder = None
        self._init_encoder()
    
    def _init_encoder(self):
        """Initialize audio encoder."""
        try:
            if self.config.audio_encoder == "whisper":
                from transformers import WhisperProcessor, WhisperModel
                self.processor = WhisperProcessor.from_pretrained("openai/whisper-base")
                self.encoder = WhisperModel.from_pretrained("openai/whisper-base")
                self.encoder_type = "whisper"
            else:
                self.encoder_type = "mfcc"
            logger.info(f"Audio encoder initialized: {self.encoder_type}")
        except ImportError:
            logger.warning("transformers not installed. Using MFCC encoding.")
            self.encoder_type = "mfcc"
        except Exception as e:
            logger.warning(f"Could not load audio encoder: {e}")
            self.encoder_type = "mfcc"
    
    def extract(self, content: bytes) -> Tuple[Optional[np.ndarray], Optional[str], Dict]:
        """
        Extract embeddings and transcription from audio.
        
        Returns:
            Tuple of (embeddings, transcription, metadata)
        """
        metadata = {"duration": 0, "sample_rate": 0}
        
        try:
            import librosa
            import soundfile as sf
            
            # Load audio
            audio, sr = sf.read(io.BytesIO(content))
            
            # Resample if needed
            if sr != self.config.audio_sample_rate:
                audio = librosa.resample(audio, orig_sr=sr, target_sr=self.config.audio_sample_rate)
                sr = self.config.audio_sample_rate
            
            metadata["duration"] = len(audio) / sr
            metadata["sample_rate"] = sr
            
            if self.encoder_type == "whisper" and self.encoder is not None:
                return self._extract_whisper(audio, metadata)
            else:
                return self._extract_mfcc(audio, sr, metadata)
        except ImportError:
            logger.warning("librosa/soundfile not installed")
            return None, None, metadata
        except Exception as e:
            logger.error(f"Error extracting audio: {e}")
            return None, None, metadata
    
    def _extract_whisper(self, audio: np.ndarray, metadata: Dict) -> Tuple[np.ndarray, str, Dict]:
        """Extract using Whisper encoder."""
        inputs = self.processor(audio, sampling_rate=self.config.audio_sample_rate, return_tensors="pt")
        
        # Get encoder outputs
        encoder_outputs = self.encoder.encoder(**inputs)
        embeddings = encoder_outputs.last_hidden_state.detach().numpy().squeeze(0)
        
        # Transcribe (simplified)
        transcription = ""  # Full transcription would require decoder
        
        metadata["encoder"] = "whisper"
        metadata["num_frames"] = embeddings.shape[0]
        
        return embeddings, transcription, metadata
    
    def _extract_mfcc(self, audio: np.ndarray, sr: int, metadata: Dict) -> Tuple[np.ndarray, None, Dict]:
        """Extract MFCC features."""
        import librosa
        
        # Compute MFCCs
        mfccs = librosa.feature.mfcc(
            y=audio,
            sr=sr,
            n_mfcc=40,
            hop_length=512,
        )
        
        # Transpose to [time, features]
        embeddings = mfccs.T.astype(np.float32)
        
        metadata["encoder"] = "mfcc"
        metadata["num_frames"] = embeddings.shape[0]
        
        return embeddings, None, metadata


class VideoExtractor:
    """Extract embeddings from video."""
    
    def __init__(self, config: TokenizationPipelineConfig):
        self.config = config
        self.image_extractor = ImageExtractor(config)
    
    def extract(self, content: bytes) -> Tuple[Optional[np.ndarray], Optional[str], Dict]:
        """
        Extract frame embeddings and audio transcription from video.
        
        Returns:
            Tuple of (frame_embeddings, audio_transcription, metadata)
        """
        metadata = {"frames": 0, "duration": 0, "fps": 0}
        
        try:
            import cv2
            import tempfile
            
            # Write to temp file (cv2 needs file path)
            with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as f:
                f.write(content)
                temp_path = f.name
            
            try:
                cap = cv2.VideoCapture(temp_path)
                
                fps = cap.get(cv2.CAP_PROP_FPS)
                total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                duration = total_frames / fps if fps > 0 else 0
                
                metadata["fps"] = fps
                metadata["duration"] = duration
                metadata["total_frames"] = total_frames
                
                # Sample frames
                target_fps = self.config.video_fps
                frame_interval = max(1, int(fps / target_fps)) if fps > 0 else 1
                max_frames = self.config.max_video_frames
                
                frame_embeddings = []
                frame_count = 0
                
                while cap.isOpened() and len(frame_embeddings) < max_frames:
                    ret, frame = cap.read()
                    if not ret:
                        break
                    
                    if frame_count % frame_interval == 0:
                        # Convert BGR to RGB
                        frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                        
                        # Encode frame as bytes
                        from PIL import Image
                        img = Image.fromarray(frame_rgb)
                        buf = io.BytesIO()
                        img.save(buf, format="PNG")
                        frame_bytes = buf.getvalue()
                        
                        # Extract embeddings
                        embeddings, _ = self.image_extractor.extract(frame_bytes)
                        if embeddings is not None:
                            # Average pool patches to single embedding per frame
                            frame_emb = embeddings.mean(axis=0)
                            frame_embeddings.append(frame_emb)
                    
                    frame_count += 1
                
                cap.release()
                
                metadata["frames"] = len(frame_embeddings)
                
                if frame_embeddings:
                    return np.stack(frame_embeddings), None, metadata
                return None, None, metadata
            finally:
                os.unlink(temp_path)
        except ImportError:
            logger.warning("OpenCV not installed")
            return None, None, metadata
        except Exception as e:
            logger.error(f"Error extracting video: {e}")
            return None, None, metadata


class ContentExtractor:
    """
    Main content extraction orchestrator.
    
    Routes content to appropriate extractors based on type.
    """
    
    def __init__(self, config: PipelineConfig):
        self.config = config
        self.tokenization_config = config.tokenization
        
        # Initialize extractors
        self.text_extractor = TextExtractor()
        self.document_extractor = DocumentExtractor()
        self.web_extractor = WebExtractor()
        self.image_extractor = ImageExtractor(self.tokenization_config)
        self.audio_extractor = AudioExtractor(self.tokenization_config)
        self.video_extractor = VideoExtractor(self.tokenization_config)
        
        # Statistics
        self.stats = {
            "total_extracted": 0,
            "by_type": {},
            "errors": 0,
        }
    
    def extract(self, item: Union[IngestedItem, PreprocessedItem]) -> ExtractedContent:
        """
        Extract content from an ingested or preprocessed item.
        
        Args:
            item: Item to extract content from
            
        Returns:
            ExtractedContent with text and/or embeddings
        """
        # Get source info
        if isinstance(item, PreprocessedItem):
            source = item.original.source
            content = item.original.content
            content_type = item.original.content_type
            
            # Use preprocessed text if available
            if item.normalized and item.content:
                return ExtractedContent(
                    source_path=source.path,
                    source_type=content_type,
                    text=item.content,
                    content_types=[ContentType.TEXT],
                    metadata=source.metadata,
                )
        else:
            source = item.source
            content = item.content
            content_type = item.content_type
        
        try:
            result = self._extract_by_type(content, content_type, source.metadata)
            result.source_path = source.path
            result.source_type = content_type
            
            self.stats["total_extracted"] += 1
            ext = source.metadata.get("extension", "unknown")
            self.stats["by_type"][ext] = self.stats["by_type"].get(ext, 0) + 1
            
            return result
        except Exception as e:
            logger.error(f"Error extracting content: {e}")
            self.stats["errors"] += 1
            
            return ExtractedContent(
                source_path=source.path,
                source_type=content_type,
                extraction_success=False,
                error_message=str(e),
            )
    
    def _extract_by_type(
        self, 
        content: Union[str, bytes], 
        content_type: str,
        metadata: Dict,
    ) -> ExtractedContent:
        """Route to appropriate extractor based on content type."""
        
        extension = metadata.get("extension", "").lower()
        
        # Text types
        if content_type.startswith("text/") or extension in [".txt", ".md", ".rst"]:
            if extension == ".md":
                text = self.text_extractor.extract_markdown(content)
            else:
                text = self.text_extractor.extract_txt(content)
            
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=text,
                content_types=[ContentType.TEXT],
                metadata=metadata,
            )
        
        # JSON
        elif content_type == "application/json" or extension == ".json":
            text, structured = self.text_extractor.extract_json(content)
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=text,
                structured_data=structured,
                content_types=[ContentType.TEXT, ContentType.STRUCTURED_DATA],
                metadata=metadata,
            )
        
        # CSV
        elif extension == ".csv":
            text, rows = self.text_extractor.extract_csv(content)
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=text,
                structured_data={"rows": rows},
                content_types=[ContentType.TEXT, ContentType.STRUCTURED_DATA],
                metadata=metadata,
            )
        
        # HTML
        elif content_type == "text/html" or extension in [".html", ".htm"]:
            text, html_meta = self.web_extractor.extract_html(content)
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=text,
                content_types=[ContentType.TEXT],
                metadata={**metadata, **html_meta},
            )
        
        # XML
        elif "xml" in content_type or extension == ".xml":
            text, xml_meta = self.web_extractor.extract_xml(content)
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=text,
                content_types=[ContentType.TEXT],
                metadata={**metadata, **xml_meta},
            )
        
        # PDF
        elif content_type == "application/pdf" or extension == ".pdf":
            if isinstance(content, str):
                content = content.encode()
            text, pdf_meta = self.document_extractor.extract_pdf(content)
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=text,
                content_types=[ContentType.TEXT],
                metadata={**metadata, **pdf_meta},
            )
        
        # DOCX
        elif extension == ".docx":
            if isinstance(content, str):
                content = content.encode()
            text, doc_meta = self.document_extractor.extract_docx(content)
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=text,
                content_types=[ContentType.TEXT],
                metadata={**metadata, **doc_meta},
            )
        
        # PPTX
        elif extension == ".pptx":
            if isinstance(content, str):
                content = content.encode()
            text, ppt_meta = self.document_extractor.extract_pptx(content)
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=text,
                content_types=[ContentType.TEXT],
                metadata={**metadata, **ppt_meta},
            )
        
        # Images
        elif content_type.startswith("image/") or extension in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"]:
            if isinstance(content, str):
                content = content.encode()
            embeddings, img_meta = self.image_extractor.extract(content)
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                image_embeddings=embeddings,
                content_types=[ContentType.IMAGE_EMBEDDING] if embeddings is not None else [],
                metadata={**metadata, **img_meta},
            )
        
        # Audio
        elif content_type.startswith("audio/") or extension in [".mp3", ".wav", ".flac", ".ogg", ".m4a"]:
            if isinstance(content, str):
                content = content.encode()
            embeddings, transcription, audio_meta = self.audio_extractor.extract(content)
            
            content_types = []
            if embeddings is not None:
                content_types.append(ContentType.AUDIO_EMBEDDING)
            if transcription:
                content_types.append(ContentType.TEXT)
            
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=transcription,
                audio_embeddings=embeddings,
                content_types=content_types,
                metadata={**metadata, **audio_meta},
            )
        
        # Video
        elif content_type.startswith("video/") or extension in [".mp4", ".avi", ".mkv", ".mov", ".webm"]:
            if isinstance(content, str):
                content = content.encode()
            embeddings, transcription, video_meta = self.video_extractor.extract(content)
            
            content_types = []
            if embeddings is not None:
                content_types.append(ContentType.VIDEO_EMBEDDING)
            if transcription:
                content_types.append(ContentType.TEXT)
            
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=transcription,
                video_embeddings=embeddings,
                content_types=content_types,
                metadata={**metadata, **video_meta},
            )
        
        # Code files
        elif extension in [".py", ".js", ".ts", ".java", ".cpp", ".c", ".h", ".go", ".rs"]:
            text = self.text_extractor.extract_txt(content)
            return ExtractedContent(
                source_path="",
                source_type=content_type,
                text=f"[CODE:{extension}]\n{text}",
                content_types=[ContentType.TEXT],
                metadata={**metadata, "is_code": True, "language": extension[1:]},
            )
        
        # Default: try as text
        else:
            try:
                if isinstance(content, bytes):
                    text = content.decode("utf-8", errors="ignore")
                else:
                    text = content
                
                return ExtractedContent(
                    source_path="",
                    source_type=content_type,
                    text=text,
                    content_types=[ContentType.TEXT],
                    metadata=metadata,
                )
            except Exception:
                return ExtractedContent(
                    source_path="",
                    source_type=content_type,
                    extraction_success=False,
                    error_message="Could not decode content",
                    metadata=metadata,
                )
    
    def get_statistics(self) -> Dict[str, Any]:
        """Get extraction statistics."""
        return dict(self.stats)
